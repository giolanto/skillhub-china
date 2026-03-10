#!/usr/bin/env python3
"""示例：自定义PPT生成"""

import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'scripts'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from themes import get_theme_colors

# ========== 配置 ==========
OUTPUT = 'custom_demo.pptx'
TITLE = '自定义演示'
SUBTITLE = 'Python生成'

# 蓝红配色
colors = get_theme_colors('blue-red')

# ========== 初始化 ==========
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== 封面 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = colors['primary']

tb = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.text = TITLE
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = colors['white']
p.alignment = PP_ALIGN.CENTER

# ========== 内容页 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = colors['white']

# 装饰条
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
shape.fill.solid()
shape.fill.fore_color.rgb = colors['secondary']
shape.line.fill.background()

# 标题
tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.9))
p = tb.text_frame.paragraphs[0]
p.text = '📌 内容页示例'
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = colors['primary']

# 正文
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
p = tb.text_frame.paragraphs[0]
p.text = '• 这是一个示例'
p.font.size = Pt(20)
p.font.color.rgb = colors['dark']

p = tb.text_frame.add_paragraph()
p.text = '• 你可以添加多行'
p.font.size = Pt(20)
p.font.color.rgb = colors['dark']

p = tb.text_frame.add_paragraph()
p.text = '• 支持多种样式'
p.font.size = Pt(20)
p.font.color.rgb = colors['dark']

# ========== 结束页 ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
tb = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.text = '感谢观看'
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = colors['primary']
p.alignment = PP_ALIGN.CENTER

# ========== 保存 ==========
prs.save(OUTPUT)
print(f'✅ 已生成: {OUTPUT}')
