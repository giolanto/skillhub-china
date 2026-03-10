#!/usr/bin/env python3
"""PPT页面模板"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def add_pic(slide, img_dir, img_name, x, y, width):
    """添加图片（图层顺序：确保最后调用）"""
    if not img_name:
        return
    path = os.path.join(img_dir, img_name) if img_dir else img_name
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(width))


def add_bar(slide, color):
    """顶部装饰条"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.12)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title(slide, text, color, y=0.3):
    """页面标题"""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = color


# ========== 页面模板 ==========

def create_cover_slide(prs, title, subtitle, colors, img_dir, bg_image):
    """封面页"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景色
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = colors['primary']
    
    # 图片背景（先放图片）
    if bg_image:
        add_pic(s, img_dir, bg_image, 0, 0, 13.333)
    
    # 半透明遮罩（再放遮罩）
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['primary']
    shape.fill.transparency = 30
    shape.line.fill.background()
    
    # 标题（最后放，最上层）
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        tb2 = s.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(1))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = colors['light_primary']
        p2.alignment = PP_ALIGN.CENTER


def create_toc_slide(prs, colors):
    """目录页"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = colors['white']
    
    add_bar(s, colors['secondary'])
    add_title(s, '📋 目录', colors['primary'])
    
    items = [
        '01 什么是...',
        '02 核心能力',
        '03 应用场景',
        '04 典型案例',
        '05 总结与展望'
    ]
    
    for i, item in enumerate(items):
        # 序号圆圈
        shape = s.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1), Inches(1.5 + i*1.1), Inches(0.6), Inches(0.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors['secondary']
        shape.line.fill.background()
        
        # 文字
        tb = s.shapes.add_textbox(Inches(1.8), Inches(1.5 + i*1.1), Inches(10), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = colors['dark']


def create_intro_slide(prs, items, colors, img_dir, image):
    """介绍页"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = colors['white']
    
    add_bar(s, colors['secondary'])
    add_title(s, '🤖 01 什么是...', colors['primary'])
    
    # 文字内容
    if items:
        tb = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(5))
        p = tb.text_frame.paragraphs[0]
        p.text = '🦞 产品简介'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = colors['secondary']
        
        for item in items:
            p = tb.text_frame.add_paragraph()
            p.text = f'• {item}'
            p.font.size = Pt(18)
            p.font.color.rgb = colors['dark']
            p.space_before = Pt(10)
    
    # 图片（最后放）
    if image:
        add_pic(s, img_dir, image, 7, 2, 5.5)


def create_features_slide(prs, features, colors, img_dir, image):
    """核心能力页（2x2卡片 + 图片）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = colors['white']
    
    add_bar(s, colors['secondary'])
    add_title(s, '💪 02 核心能力', colors['primary'])
    
    # 2x2卡片布局
    for i, item in enumerate(features[:4]):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 3.2
        y = 1.3 + row * 2.8
        
        icon = item.get('icon', '')
        title = item.get('title', '')
        desc = item.get('desc', '')
        
        # 卡片背景
        shape = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(3), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors['light_primary']
        shape.line.fill.background()
        
        # 标题
        tb = s.shapes.add_textbox(Inches(x+0.2), Inches(y+0.3), Inches(2.6), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = f'{icon} {title}'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors['primary']
        
        # 描述
        tb2 = s.shapes.add_textbox(Inches(x+0.2), Inches(y+1.0), Inches(2.6), Inches(1))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = colors['dark']
    
    # 图片（最后放）
    if image:
        add_pic(s, img_dir, image, 7, 1.5, 6)


def create_scenarios_slide(prs, scenarios, colors, img_dir, image):
    """应用场景页（左侧列表 + 右侧图片）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = colors['white']
    
    add_bar(s, colors['secondary'])
    add_title(s, '📱 03 应用场景', colors['primary'])
    
    # 左侧场景
    for i, item in enumerate(scenarios[:4]):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 3.5
        y = 1.3 + row * 3.0
        
        icon = item.get('icon', '')
        title = item.get('title', '')
        desc = item.get('desc', '')
        
        # 红色竖条
        shape = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(0.1), Inches(2.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors['secondary']
        shape.line.fill.background()
        
        # 标题
        tb = s.shapes.add_textbox(Inches(x+0.3), Inches(y+0.2), Inches(3), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = f'{icon} {title}'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors['primary']
        
        # 描述
        tb2 = s.shapes.add_textbox(Inches(x+0.3), Inches(y+0.7), Inches(3), Inches(1.5))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = colors['dark']
    
    # 图片（最后放）
    if image:
        add_pic(s, img_dir, image, 8, 1.5, 5)


def create_cases_slide(prs, cases, colors):
    """典型案例页"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = colors['white']
    
    add_bar(s, colors['secondary'])
    add_title(s, '📌 04 典型案例', colors['primary'])
    
    # 2x2卡片
    for i, item in enumerate(cases[:4]):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 6.5
        y = 1.3 + row * 2.8
        
        icon = item.get('icon', '')
        title = item.get('title', '')
        desc = item.get('desc', '')
        
        # 卡片
        shape = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(6.2), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors['light_secondary']
        shape.line.fill.background()
        
        # 标题
        tb = s.shapes.add_textbox(Inches(x+0.3), Inches(y+0.6), Inches(5.6), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = f'{icon} {title}'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = colors['secondary']


def create_summary_slide(prs, items, colors):
    """总结页"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = colors['white']
    
    add_bar(s, colors['secondary'])
    add_title(s, '🌟 05 总结与展望', colors['primary'])
    
    # 2x2布局
    for i, item in enumerate(items[:4]):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 6.5
        y = 1.4 + row * 2.8
        
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(6.2), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = colors['secondary']


def create_ending_slide(prs, colors, img_dir, bg_image):
    """结束页"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 图片背景（先放）
    if bg_image:
        add_pic(s, img_dir, bg_image, 0, 0, 13.333)
    
    # 文字（后放，最上层）
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = '感谢关注'
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    p.alignment = PP_ALIGN.CENTER
    
    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(1))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = '让AI成为你的左膀右臂'
    p2.font.size = Pt(28)
    p2.font.color.rgb = colors['dark']
    p2.alignment = PP_ALIGN.CENTER
